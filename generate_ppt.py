"""
generate_ppt.py  –  AI Bus Ticketing System PowerPoint Generator
Run:  python generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

# ── Colour palette ─────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0F, 0x34, 0x60)   # #0f3460
PURPLE      = RGBColor(0x6C, 0x63, 0xFF)   # #6c63ff
TEAL        = RGBColor(0x00, 0xD4, 0xAA)   # #00d4aa
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG    = RGBColor(0xF0, 0xF4, 0xFF)   # #f0f4ff
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)   # #1a1a2e
MUTED       = RGBColor(0x55, 0x55, 0x77)   # #555577
GOLD        = RGBColor(0xF5, 0xA6, 0x23)   # #f5a623
RED_ACCENT  = RGBColor(0xE9, 0x45, 0x60)   # #e94560
SLIDE_BG    = RGBColor(0x12, 0x12, 0x2E)   # near-black navy

W = Inches(13.33)   # Widescreen 16:9
H = Inches(7.5)

OUT = "AI_Bus_Ticketing_System_Presentation.pptx"


# ── Helpers ────────────────────────────────────────────────────────────────────

def prs_init():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]   # completely blank


def add_rect(slide, left, top, width, height, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def slide_bg(slide, color=SLIDE_BG):
    """Fill slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def accent_bar(slide, top=Inches(0), height=Inches(0.08), color=PURPLE):
    add_rect(slide, 0, top, W, height, color)


def top_strip(slide):
    """Thin coloured top bar for non-title slides."""
    add_rect(slide, 0, 0, W, Inches(0.12), PURPLE)


def bottom_strip(slide, slide_num=None, total=None):
    add_rect(slide, 0, H - Inches(0.45), W, Inches(0.45), NAVY)
    add_text_box(slide,
        "AI Chat Bot Based Bus Ticketing System",
        Inches(0.3), H - Inches(0.42), Inches(7), Inches(0.4),
        font_size=9, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.LEFT
    )
    if slide_num and total:
        add_text_box(slide,
            f"{slide_num} / {total}",
            W - Inches(1.2), H - Inches(0.42), Inches(1), Inches(0.4),
            font_size=9, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.RIGHT
        )


def section_label(slide, text, top=Inches(0.18), color=TEAL):
    add_text_box(slide, text,
        Inches(0.5), top, Inches(12), Inches(0.4),
        font_size=11, bold=True, color=color, align=PP_ALIGN.LEFT
    )


def slide_title(slide, text, top=Inches(0.55), size=32, color=WHITE):
    add_text_box(slide, text,
        Inches(0.5), top, Inches(12.3), Inches(0.9),
        font_size=size, bold=True, color=color, align=PP_ALIGN.LEFT
    )


def divider(slide, top, color=PURPLE, width_frac=0.12):
    add_rect(slide, Inches(0.5), top, Inches(width_frac * 13.33), Inches(0.04), color)


# ── Bullet helper ──────────────────────────────────────────────────────────────

def bullet_box(slide, bullets, left, top, width, height,
               font_size=14, title=None, title_color=TEAL,
               bullet_color=WHITE, bg=None, border_color=None):
    if bg:
        r = add_rect(slide, left, top, width, height, bg)
        if border_color:
            r.line.color.rgb = border_color
            r.line.width = Pt(1)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.12),
        width - Inches(0.3), height - Inches(0.24)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = title_color
    for i, b in enumerate(bullets):
        para = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = ("  " + b) if title else b
        run.font.size = Pt(font_size)
        run.font.color.rgb = bullet_color
        run.font.bold = False


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title slide."""
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, SLIDE_BG)

    # Large purple gradient-like left panel
    add_rect(slide, 0, 0, Inches(5.5), H, NAVY)
    add_rect(slide, 0, 0, Inches(0.18), H, PURPLE)

    # Bus icon placeholder text
    add_text_box(slide, "BUS", Inches(0.3), Inches(1.2), Inches(5), Inches(1.8),
                 font_size=90, bold=True, color=RGBColor(0x1F, 0x44, 0x70), align=PP_ALIGN.CENTER)

    # Title on right
    add_text_box(slide, "AI CHAT BOT BASED",
                 Inches(5.8), Inches(1.0), Inches(7.2), Inches(0.9),
                 font_size=30, bold=True, color=PURPLE, align=PP_ALIGN.LEFT)
    add_text_box(slide, "BUS TICKETING SYSTEM",
                 Inches(5.8), Inches(1.85), Inches(7.2), Inches(0.9),
                 font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_rect(slide, Inches(5.8), Inches(2.9), Inches(4.5), Inches(0.06), TEAL)

    add_text_box(slide,
        "A full-stack AI-powered travel assistant that finds,\n"
        "filters and presents live bus tickets through\n"
        "natural-language conversation.",
        Inches(5.8), Inches(3.1), Inches(7.1), Inches(1.4),
        font_size=15, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.LEFT
    )

    # Stat chips
    stats = [("3+","AI Models"), ("10","Bus Scrapers"), ("15+","Cities"), ("4","Smart Filters")]
    for i, (num, lbl) in enumerate(stats):
        x = Inches(5.8 + i * 1.75)
        add_rect(slide, x, Inches(4.7), Inches(1.6), Inches(0.9), RGBColor(0x1E, 0x2A, 0x50))
        add_text_box(slide, num, x, Inches(4.72), Inches(1.6), Inches(0.45),
                     font_size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        add_text_box(slide, lbl, x, Inches(5.1), Inches(1.6), Inches(0.4),
                     font_size=9, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)

    add_text_box(slide,
        f"Presentation  |  {datetime.datetime.now().strftime('%B %Y')}",
        Inches(5.8), Inches(6.8), Inches(7), Inches(0.4),
        font_size=10, color=MUTED, align=PP_ALIGN.LEFT
    )
    bottom_strip(slide)


def slide_02_agenda(prs, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    top_strip(slide)
    bottom_strip(slide, 2, total)
    section_label(slide, "OVERVIEW")
    slide_title(slide, "Agenda")
    divider(slide, Inches(1.35))

    items = [
        ("01", "Project Overview & Problem Statement"),
        ("02", "System Architecture"),
        ("03", "Technology Stack  (Frontend)"),
        ("04", "Technology Stack  (Backend & AI)"),
        ("05", "AI Intent Pipeline"),
        ("06", "Smart Filtering & Sorting"),
        ("07", "Live Data Scraping"),
        ("08", "Core UI Features"),
        ("09", "Advanced Features"),
        ("10", "Data Flow Diagram"),
        ("11", "Limitations & Future Scope"),
        ("12", "Conclusion"),
    ]

    cols = [items[:6], items[6:]]
    for ci, col in enumerate(cols):
        for ri, (num, label) in enumerate(col):
            x = Inches(0.5 + ci * 6.5)
            y = Inches(1.6 + ri * 0.85)
            add_rect(slide, x, y, Inches(0.5), Inches(0.55), PURPLE)
            add_text_box(slide, num, x, y, Inches(0.5), Inches(0.55),
                         font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_text_box(slide, label, x + Inches(0.6), y + Inches(0.08),
                         Inches(5.6), Inches(0.45), font_size=14, color=WHITE)


def slide_03_overview(prs, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    top_strip(slide)
    bottom_strip(slide, 3, total)
    section_label(slide, "SLIDE 01")
    slide_title(slide, "Project Overview & Problem Statement")
    divider(slide, Inches(1.35))

    # Problem box
    bullet_box(slide,
        [
            "Bus booking in India is fragmented across RedBus, AbhiBus, MakeMyTrip, etc.",
            "Users must manually filter by price, seat type, operator, timing.",
            "Non-technical travellers find complex portals confusing.",
            "No single unified interface that understands plain English.",
        ],
        Inches(0.5), Inches(1.55), Inches(5.9), Inches(2.5),
        font_size=13, title="The Problem",
        title_color=RED_ACCENT, bullet_color=RGBColor(0xDD, 0xDD, 0xFF),
        bg=RGBColor(0x1A, 0x1E, 0x3A), border_color=RED_ACCENT
    )

    # Solution box
    bullet_box(slide,
        [
            "A smart AI chatbot that accepts plain English queries.",
            "Extracts travel intent, filters, dates automatically.",
            "Fetches live data from RedBus via Playwright scraping.",
            "Presents clean, filtered, sorted results instantly.",
        ],
        Inches(6.7), Inches(1.55), Inches(5.9), Inches(2.5),
        font_size=13, title="Our Solution",
        title_color=TEAL, bullet_color=RGBColor(0xDD, 0xDD, 0xFF),
        bg=RGBColor(0x1A, 0x1E, 0x3A), border_color=TEAL
    )

    # Quote / tagline
    add_rect(slide, Inches(0.5), Inches(4.3), Inches(12.2), Inches(1.15),
             RGBColor(0x0F, 0x34, 0x60))
    add_text_box(slide,
        '"Find me the cheapest AC bus from Madurai to Chennai tomorrow"',
        Inches(0.7), Inches(4.38), Inches(11.8), Inches(0.55),
        font_size=17, bold=True, italic=True, color=GOLD, align=PP_ALIGN.CENTER
    )
    add_text_box(slide,
        "— The system understands this, fetches live data, filters, and shows results in seconds.",
        Inches(0.7), Inches(4.88), Inches(11.8), Inches(0.45),
        font_size=12, color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.CENTER
    )

    # Key goals row
    goals = [
        ("Natural Language\nInterface", PURPLE),
        ("Live Real-Time\nBus Data", TEAL),
        ("Smart AI\nFiltering", GOLD),
        ("Premium\nUI/UX", RED_ACCENT),
    ]
    for i, (g, c) in enumerate(goals):
        x = Inches(0.5 + i * 3.1)
        add_rect(slide, x, Inches(5.7), Inches(2.9), Inches(0.9), RGBColor(0x1E, 0x28, 0x50))
        add_rect(slide, x, Inches(5.7), Inches(0.08), Inches(0.9), c)
        add_text_box(slide, g, x + Inches(0.2), Inches(5.72), Inches(2.6), Inches(0.86),
                     font_size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def slide_04_architecture(prs, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    top_strip(slide)
    bottom_strip(slide, 4, total)
    section_label(slide, "SLIDE 02")
    slide_title(slide, "System Architecture")
    divider(slide, Inches(1.35))

    layers = [
        ("LAYER 1 — Frontend",   "React 18 + TypeScript + Vite",  "User Interface, Chat, Results, Map, Theme Toggle", PURPLE),
        ("LAYER 2 — Backend",    "FastAPI + Python 3.11+",         "Intent Parsing, API Routing, Filter & Sort Logic",  TEAL),
        ("LAYER 3 — AI Engine",  "OpenRouter (LLM) + Rule Parser", "Natural Language Understanding, Intent Extraction",  GOLD),
        ("LAYER 4 — Scraper",    "Playwright (async)",              "Live Bus Data from RedBus GraphQL API",            RED_ACCENT),
        ("LAYER 5 — Storage",    "localStorage + .env",             "Theme, Sessions, Booking History, API Keys",       RGBColor(0xCC, 0x77, 0xFF)),
    ]

    for i, (layer, tech, desc, color) in enumerate(layers):
        y = Inches(1.55 + i * 1.02)
        add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.88), RGBColor(0x16, 0x1E, 0x3A))
        add_rect(slide, Inches(0.5), y, Inches(0.1), Inches(0.88), color)
        add_text_box(slide, layer, Inches(0.75), y + Inches(0.06),
                     Inches(3.0), Inches(0.38), font_size=11, bold=True, color=color)
        add_text_box(slide, tech, Inches(0.75), y + Inches(0.44),
                     Inches(3.0), Inches(0.35), font_size=10, color=RGBColor(0xCC, 0xCC, 0xFF))
        # Arrow
        add_rect(slide, Inches(3.9), y + Inches(0.3), Inches(0.5), Inches(0.06),
                 RGBColor(0x44, 0x44, 0x77))
        add_text_box(slide, desc, Inches(4.5), y + Inches(0.22),
                     Inches(8.2), Inches(0.55), font_size=13, color=WHITE)


def slide_05_tech_frontend(prs, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    top_strip(slide)
    bottom_strip(slide, 5, total)
    section_label(slide, "SLIDE 03")
    slide_title(slide, "Technology Stack — Frontend")
    divider(slide, Inches(1.35))

    cards = [
        ("React 18",       "Component-based UI framework.\nHandles state, routing & rendering.",    PURPLE),
        ("TypeScript 5",   "Strongly-typed JavaScript.\nCatch bugs at compile time.",               TEAL),
        ("Vite 6",         "Ultra-fast build tool &\ndevelopment server (HMR).",                    GOLD),
        ("Vanilla CSS",    "No Tailwind — full CSS control.\nGlassmorphism, variables, animations.",RED_ACCENT),
        ("Leaflet.js",     "Interactive map for tourist spots\nnear the destination city.",          RGBColor(0x4C, 0xAF, 0x50)),
        ("Lucide React",   "Clean icon library for\nUI elements and badges.",                       RGBColor(0x03, 0xA9, 0xF4)),
        ("Web Speech API", "Voice input support for\nhands-free bus search.",                       RGBColor(0xFF, 0x57, 0x22)),
        ("localStorage",   "Persist theme, history,\nchat sessions client-side.",                   RGBColor(0xFF, 0xC1, 0x07)),
    ]

    cols = 4
    for i, (name, desc, color) in enumerate(cards):
        row = i // cols
        col = i % cols
        x = Inches(0.4 + col * 3.15)
        y = Inches(1.6 + row * 2.3)
        add_rect(slide, x, y, Inches(2.9), Inches(2.0), RGBColor(0x16, 0x1E, 0x3A))
        add_rect(slide, x, y, Inches(2.9), Inches(0.1), color)
        add_text_box(slide, name, x + Inches(0.1), y + Inches(0.2),
                     Inches(2.7), Inches(0.45), font_size=14, bold=True, color=color)
        add_text_box(slide, desc, x + Inches(0.1), y + Inches(0.65),
                     Inches(2.7), Inches(1.15), font_size=11, color=RGBColor(0xCC, 0xCC, 0xFF))


def slide_06_tech_backend(prs, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    top_strip(slide)
    bottom_strip(slide, 6, total)
    section_label(slide, "SLIDE 04")
    slide_title(slide, "Technology Stack — Backend & AI")
    divider(slide, Inches(1.35))

    left_items = [
        ("Python 3.11+",   "Core backend language", PURPLE),
        ("FastAPI",         "REST API framework — async, fast, auto-docs", TEAL),
        ("Uvicorn",         "ASGI server running the FastAPI app", GOLD),
        ("Pydantic",        "Data validation & request schemas", RED_ACCENT),
        ("httpx",           "Async HTTP client for AI API calls", RGBColor(0x03, 0xA9, 0xF4)),
        ("python-dotenv",   "Load API keys from .env file", RGBColor(0x4C, 0xAF, 0x50)),
        ("Docker",          "Containerisation for easy deployment", RGBColor(0xFF, 0x57, 0x22)),
    ]

    right_items = [
        ("OpenRouter API",  "Gateway to 100+ free/paid LLMs", PURPLE),
        ("Nous Hermes 3\nLlama 3.1 405B", "Primary AI model for intent parsing", TEAL),
        ("Rule-Based Parser","Regex fallback when AI hits rate limit", GOLD),
        ("Playwright",      "Async browser automation for scraping", RED_ACCENT),
        ("RedBus GraphQL",  "Live bus data interception via API", RGBColor(0xFF, 0xC1, 0x07)),
        ("BeautifulSoup",   "HTML parsing fallback if needed", RGBColor(0x4C, 0xAF, 0x50)),
        ("Git",             "Version control throughout development", RGBColor(0xAA, 0xAA, 0xFF)),
    ]

    for i, (name, desc, color) in enumerate(left_items):
        y = Inches(1.6 + i * 0.74)
        add_rect(slide, Inches(0.4), y, Inches(6.1), Inches(0.65), RGBColor(0x16, 0x1E, 0x3A))
        add_rect(slide, Inches(0.4), y, Inches(0.08), Inches(0.65), color)
        add_text_box(slide, name, Inches(0.6), y + Inches(0.04),
                     Inches(2.2), Inches(0.3), font_size=12, bold=True, color=color)
        add_text_box(slide, desc, Inches(2.9), y + Inches(0.14),
                     Inches(3.5), Inches(0.35), font_size=11, color=WHITE)

    for i, (name, desc, color) in enumerate(right_items):
        y = Inches(1.6 + i * 0.74)
        add_rect(slide, Inches(6.8), y, Inches(6.1), Inches(0.65), RGBColor(0x16, 0x1E, 0x3A))
        add_rect(slide, Inches(6.8), y, Inches(0.08), Inches(0.65), color)
        add_text_box(slide, name, Inches(7.0), y + Inches(0.04),
                     Inches(2.2), Inches(0.55), font_size=12, bold=True, color=color)
        add_text_box(slide, desc, Inches(9.3), y + Inches(0.14),
                     Inches(3.5), Inches(0.35), font_size=11, color=WHITE)


def slide_07_ai_pipeline(prs, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    top_strip(slide)
    bottom_strip(slide, 7, total)
    section_label(slide, "SLIDE 05")
    slide_title(slide, "AI Intent Pipeline")
    divider(slide, Inches(1.35))

    steps = [
        ("1", "User Input",        "User types or speaks a query like\n'AC bus from Madurai to Chennai tomorrow'", PURPLE),
        ("2", "Context Merge",     "Previous conversation context (cities,\ndates, filters) merged with new query",   TEAL),
        ("3", "AI Model Call",     "Nous Hermes 3 (via OpenRouter) returns\nstructured JSON: {from, to, date, filter, sort_by}", GOLD),
        ("4", "Rate-Limit Fallback","On HTTP 429 or empty response, tries\nnext model in chain automatically",          RED_ACCENT),
        ("5", "Rule-Based Parser", "Final fallback: regex extracts cities,\ndates, filter keywords, sort preferences",  RGBColor(0xCC, 0x77, 0xFF)),
        ("6", "Ready to Search?",  "If from/to/date all collected → scrape.\nElse → ask follow-up with quick replies",   RGBColor(0x03, 0xA9, 0xF4)),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        row = i // 3
        col = i % 3
        x = Inches(0.4 + col * 4.3)
        y = Inches(1.6 + row * 2.5)
        add_rect(slide, x, y, Inches(4.0), Inches(2.2), RGBColor(0x16, 0x1E, 0x3A))
        add_rect(slide, x, y, Inches(4.0), Inches(0.1), color)
        # Number badge
        add_rect(slide, x + Inches(0.15), y + Inches(0.2),
                 Inches(0.45), Inches(0.45), color)
        add_text_box(slide, num, x + Inches(0.15), y + Inches(0.2),
                     Inches(0.45), Inches(0.45), font_size=16, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, title, x + Inches(0.7), y + Inches(0.25),
                     Inches(3.1), Inches(0.4), font_size=13, bold=True, color=color)
        add_text_box(slide, desc, x + Inches(0.15), y + Inches(0.75),
                     Inches(3.7), Inches(1.3), font_size=11,
                     color=RGBColor(0xCC, 0xCC, 0xFF))

        # Arrow between cols
        if col < 2:
            add_text_box(slide, "→", x + Inches(4.05), y + Inches(0.9),
                         Inches(0.22), Inches(0.4), font_size=18, bold=True,
                         color=PURPLE, align=PP_ALIGN.CENTER)


def slide_08_filtering(prs, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    top_strip(slide)
    bottom_strip(slide, 8, total)
    section_label(slide, "SLIDE 06")
    slide_title(slide, "Smart AI Filtering & Sorting")
    divider(slide, Inches(1.35))

    filters = [
        ("cheapest", "💰 CHEAPEST",   "Sorts all results by price ascending.\nCheapest bus always shown first.",     GOLD),
        ("ac",       "❄️ AC BUS",     "Filters to show only AC-equipped buses.\n'BEST AC BUS' badge on top result.",  TEAL),
        ("sleeper",  "🛏 SLEEPER",    "Filters to sleeper-only buses.\n'BEST SLEEPER' badge shown.",                  PURPLE),
        ("non_ac",   "🚌 NON-AC",     "Shows economy non-AC buses.\nGreat for budget travellers.",                   RGBColor(0x4C, 0xAF, 0x50)),
        ("volvo",    "✨ VOLVO",       "Filters luxury/Volvo buses.\nPremium operator highlight.",                    RGBColor(0xFF, 0x57, 0x22)),
        ("fastest",  "⚡ FASTEST",    "Sorts by journey duration.\nShortest trip time shown first.",                  RGBColor(0x03, 0xA9, 0xF4)),
        ("night",    "🌙 NIGHT BUS",  "Shows buses departing after 8pm.\nIdeal for overnight travel.",               RGBColor(0xCC, 0x77, 0xFF)),
    ]

    for i, (key, label, desc, color) in enumerate(filters):
        row = i // 4
        col = i % 4
        if i == 4:  # handle 7th item centered
            pass
        x = Inches(0.4 + col * 3.2)
        if row == 1:
            x = Inches(1.95 + (i - 4) * 3.2)
        y = Inches(1.6 + row * 2.45)
        add_rect(slide, x, y, Inches(2.95), Inches(2.15), RGBColor(0x16, 0x1E, 0x3A))
        add_rect(slide, x, y, Inches(2.95), Inches(0.09), color)
        add_text_box(slide, label, x + Inches(0.12), y + Inches(0.18),
                     Inches(2.7), Inches(0.4), font_size=13, bold=True, color=color)
        add_text_box(slide, desc, x + Inches(0.12), y + Inches(0.62),
                     Inches(2.7), Inches(1.3), font_size=11,
                     color=RGBColor(0xCC, 0xCC, 0xFF))


def slide_09_scraping(prs, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    top_strip(slide)
    bottom_strip(slide, 9, total)
    section_label(slide, "SLIDE 07")
    slide_title(slide, "Live Data Scraping — How It Works")
    divider(slide, Inches(1.35))

    # Left — flow
    steps_left = [
        ("Step 1", "Playwright launches headless Chromium browser", PURPLE),
        ("Step 2", "Navigates to RedBus search URL with route + date", TEAL),
        ("Step 3", "Intercepts the internal GraphQL API response", GOLD),
        ("Step 4", "Parses JSON: operator, price, seats, times, URLs", RED_ACCENT),
        ("Step 5", "Returns up to 10 real buses per query", RGBColor(0x4C, 0xAF, 0x50)),
        ("Step 6", "Results cached 10 sec to prevent duplicate calls", RGBColor(0xCC, 0x77, 0xFF)),
    ]
    for i, (step, desc, color) in enumerate(steps_left):
        y = Inches(1.6 + i * 0.84)
        add_rect(slide, Inches(0.4), y, Inches(7.5), Inches(0.75), RGBColor(0x16, 0x1E, 0x3A))
        add_rect(slide, Inches(0.4), y, Inches(0.08), Inches(0.75), color)
        add_text_box(slide, step, Inches(0.6), y + Inches(0.08),
                     Inches(1.1), Inches(0.3), font_size=10, bold=True, color=color)
        add_text_box(slide, desc, Inches(1.75), y + Inches(0.18),
                     Inches(5.9), Inches(0.4), font_size=12, color=WHITE)

    # Right — data returned
    bullet_box(slide,
        [
            "  Bus operator name (e.g. SETC, FlixBus)",
            "  Departure & arrival time",
            "  Journey duration",
            "  Price (in INR)",
            "  Seats available (+ urgency flag if < 5)",
            "  Bus type (AC / Sleeper / Volvo / etc.)",
            "  Amenities (WiFi, charging, live tracking)",
            "  Direct booking URL to RedBus",
            "  Free cancellation flag",
        ],
        Inches(8.2), Inches(1.6), Inches(4.8), Inches(5.2),
        font_size=12, title="Data Returned Per Bus",
        title_color=TEAL, bullet_color=RGBColor(0xDD, 0xDD, 0xFF),
        bg=RGBColor(0x14, 0x1E, 0x3A), border_color=TEAL
    )


def slide_10_ui_features(prs, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    top_strip(slide)
    bottom_strip(slide, 10, total)
    section_label(slide, "SLIDE 08")
    slide_title(slide, "Core UI Features")
    divider(slide, Inches(1.35))

    features = [
        ("Chat Interface",
         "Multi-turn conversation with message history,\nquick-reply chips, voice input & bot responses.", PURPLE),
        ("Bus Results Panel",
         "Ticket cards with operator logo, price, times,\nseats, badges, save/expand/book actions.", TEAL),
        ("Smart Filter Badges",
         "Active filter banner (AC / Cheapest / Sleeper)\nwith clear button. Badge on top result card.", GOLD),
        ("Dark / Light Theme",
         "2026 glassmorphism design. Theme toggle button\nsaves preference in localStorage.", RED_ACCENT),
        ("Connecting Routes",
         "Auto-suggests 2-leg journey when no direct bus\nexists (e.g. Chennai → Dindigul → Kodaikanal).", RGBColor(0x4C, 0xAF, 0x50)),
        ("Booking History",
         "In-browser booking tracker with confirmation IDs,\noperator, date, seat, amount stored.", RGBColor(0x03, 0xA9, 0xF4)),
    ]

    for i, (name, desc, color) in enumerate(features):
        row = i // 3
        col = i % 3
        x = Inches(0.4 + col * 4.3)
        y = Inches(1.6 + row * 2.5)
        add_rect(slide, x, y, Inches(4.0), Inches(2.2), RGBColor(0x16, 0x1E, 0x3A))
        add_rect(slide, x, y, Inches(0.08), Inches(2.2), color)
        add_text_box(slide, name, x + Inches(0.2), y + Inches(0.2),
                     Inches(3.6), Inches(0.45), font_size=14, bold=True, color=color)
        add_text_box(slide, desc, x + Inches(0.2), y + Inches(0.72),
                     Inches(3.6), Inches(1.3), font_size=11,
                     color=RGBColor(0xCC, 0xCC, 0xFF))


def slide_11_advanced(prs, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    top_strip(slide)
    bottom_strip(slide, 11, total)
    section_label(slide, "SLIDE 09")
    slide_title(slide, "Advanced Features")
    divider(slide, Inches(1.35))

    advanced = [
        ("Travel Guide Panel",
         "AI generates tips, top places to visit, best food, best time.\nFor any destination city searched by the user.", PURPLE),
        ("Interactive Map",
         "Leaflet.js renders a map pinpointing tourist spots\nnear the destination. Clickable markers with names.", TEAL),
        ("Voice Input",
         "Web Speech API lets users speak their search query.\nTranscript auto-fills the chat input box.", GOLD),
        ("User Personalisation",
         "App greets user by name. Remembers frequent routes.\nSuggests 'AC Bus' / 'Tomorrow' based on past use.", RED_ACCENT),
        ("Review System",
         "Users can rate bus operators with star ratings.\nRatings displayed on ticket cards.", RGBColor(0x4C, 0xAF, 0x50)),
        ("Auth Modal",
         "Simple login with username. Session persisted.\nPersonalised greeting on return visits.", RGBColor(0xCC, 0x77, 0xFF)),
        ("Saved Tickets",
         "Users can bookmark tickets for comparison.\nSaved tab shows all favourited results.", RGBColor(0x03, 0xA9, 0xF4)),
        ("Multi-turn Dialogue",
         "Bot asks follow-up questions if info is missing.\nContext carried across conversation turns.", RGBColor(0xFF, 0x57, 0x22)),
    ]

    for i, (name, desc, color) in enumerate(advanced):
        row = i // 4
        col = i % 4
        x = Inches(0.4 + col * 3.2)
        y = Inches(1.6 + row * 2.5)
        add_rect(slide, x, y, Inches(2.95), Inches(2.2), RGBColor(0x16, 0x1E, 0x3A))
        add_rect(slide, x, y, Inches(2.95), Inches(0.09), color)
        add_text_box(slide, name, x + Inches(0.12), y + Inches(0.2),
                     Inches(2.7), Inches(0.4), font_size=12, bold=True, color=color)
        add_text_box(slide, desc, x + Inches(0.12), y + Inches(0.65),
                     Inches(2.7), Inches(1.3), font_size=10.5,
                     color=RGBColor(0xCC, 0xCC, 0xFF))


def slide_12_dataflow(prs, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    top_strip(slide)
    bottom_strip(slide, 12, total)
    section_label(slide, "SLIDE 10")
    slide_title(slide, "End-to-End Data Flow")
    divider(slide, Inches(1.35))

    nodes = [
        ("User sends\nchat message",            Inches(0.5),  Inches(2.8), PURPLE),
        ("React UI\nPOST /search",              Inches(2.1),  Inches(2.8), TEAL),
        ("FastAPI\nBackend",                    Inches(3.9),  Inches(2.8), GOLD),
        ("AI / Rule\nParser",                   Inches(5.7),  Inches(2.8), RED_ACCENT),
        ("Playwright\nScraper",                 Inches(7.5),  Inches(2.8), RGBColor(0x4C, 0xAF, 0x50)),
        ("Filter &\nSort Logic",                Inches(9.3),  Inches(2.8), RGBColor(0x03, 0xA9, 0xF4)),
        ("JSON Response\nto UI",                Inches(11.1), Inches(2.8), PURPLE),
    ]
    node_w = Inches(1.5)
    node_h = Inches(0.9)

    for (text, x, y, color) in nodes:
        add_rect(slide, x, y, node_w, node_h, RGBColor(0x16, 0x1E, 0x3A))
        add_rect(slide, x, y, node_w, Inches(0.09), color)
        add_text_box(slide, text, x + Inches(0.05), y + Inches(0.15),
                     node_w - Inches(0.1), node_h - Inches(0.18),
                     font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrows between nodes
    for i in range(len(nodes) - 1):
        ax = nodes[i][1] + node_w + Inches(0.03)
        ay = nodes[i][2] + node_h / 2 - Inches(0.03)
        add_rect(slide, ax, ay, Inches(0.07), Inches(0.06), PURPLE)

    # Labels below arrows
    arrow_labels = [
        "HTTP POST", "Route &\nQuery", "NL Query", "Intent JSON", "Bus List", "Filtered\nResults"
    ]
    for i, lbl in enumerate(arrow_labels):
        ax = nodes[i][1] + node_w + Inches(0.08)
        add_text_box(slide, lbl, ax, Inches(3.85), Inches(1.4), Inches(0.55),
                     font_size=8.5, color=MUTED, align=PP_ALIGN.CENTER)

    # Bottom detail boxes
    details = [
        ("React UI", "Chat, Results, Badges,\nTheme Toggle, Map, History", PURPLE),
        ("FastAPI Backend", "CORS, Pydantic, Caching,\nConnecting Routes, Travel Guide API", TEAL),
        ("AI Engine", "OpenRouter → Hermes 3\n→ Fallback models → Rule parser", GOLD),
        ("Scraper", "Playwright intercepts\nRedBus GraphQL API", RED_ACCENT),
    ]
    for i, (name, desc, color) in enumerate(details):
        x = Inches(0.4 + i * 3.2)
        y = Inches(5.1)
        add_rect(slide, x, y, Inches(3.0), Inches(1.5), RGBColor(0x14, 0x1E, 0x3A))
        add_rect(slide, x, y, Inches(3.0), Inches(0.07), color)
        add_text_box(slide, name, x + Inches(0.1), y + Inches(0.12),
                     Inches(2.8), Inches(0.35), font_size=11, bold=True, color=color)
        add_text_box(slide, desc, x + Inches(0.1), y + Inches(0.5),
                     Inches(2.8), Inches(0.9), font_size=10, color=RGBColor(0xCC, 0xCC, 0xFF))


def slide_13_limitations(prs, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    top_strip(slide)
    bottom_strip(slide, 13, total)
    section_label(slide, "SLIDE 11")
    slide_title(slide, "Limitations & Future Scope")
    divider(slide, Inches(1.35))

    bullet_box(slide,
        [
            "Free AI tier (OpenRouter) has rate limits — rule-based fallback mitigates this",
            "Scraper depends on RedBus API structure — may need updates on site changes",
            "No real payment gateway — booking redirects to third-party sites",
            "Bus data limited to routes available on RedBus India",
            "No persistent database — all history stored in localStorage",
        ],
        Inches(0.4), Inches(1.6), Inches(6.0), Inches(4.5),
        font_size=13, title="Current Limitations",
        title_color=RED_ACCENT, bullet_color=RGBColor(0xDD, 0xDD, 0xFF),
        bg=RGBColor(0x16, 0x1E, 0x3A), border_color=RED_ACCENT
    )

    bullet_box(slide,
        [
            "Fine-tuned NLP model for richer intent understanding",
            "Add train, flight, and cab booking support",
            "Real-time seat-map & seat selection inside the app",
            "Razorpay / Stripe payment gateway integration",
            "Progressive Web App (PWA) with push notifications",
            "Tamil, Hindi, Kannada multi-language support",
            "PostgreSQL backend for persistent user data",
            "ML price prediction using historical pricing data",
        ],
        Inches(6.7), Inches(1.6), Inches(6.2), Inches(4.5),
        font_size=13, title="Future Enhancements",
        title_color=TEAL, bullet_color=RGBColor(0xDD, 0xDD, 0xFF),
        bg=RGBColor(0x16, 0x1E, 0x3A), border_color=TEAL
    )


def slide_14_conclusion(prs, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide, NAVY)
    top_strip(slide)
    bottom_strip(slide, 14, total)

    add_rect(slide, 0, 0, Inches(0.2), H, PURPLE)

    add_text_box(slide, "CONCLUSION", Inches(0.5), Inches(0.9),
                 Inches(12), Inches(0.5), font_size=11, bold=True, color=TEAL)
    add_text_box(slide,
        "What We Built",
        Inches(0.5), Inches(1.35), Inches(12), Inches(0.8),
        font_size=30, bold=True, color=WHITE)
    add_rect(slide, Inches(0.5), Inches(2.1), Inches(2.0), Inches(0.06), PURPLE)

    summary = (
        "The AI Chat Bot Based Bus Ticketing System successfully combines Conversational AI, "
        "real-time web scraping, and premium UI design to transform a fragmented bus-booking "
        "experience into a simple, zero-learning-curve chat interaction.\n\n"
        "By using a layered AI strategy — primary LLM → fallback models → rule-based parser — "
        "the system delivers near-100% reliability for intent extraction with zero inference cost.\n\n"
        "The result: a production-ready AI travel assistant covering 15+ South Indian cities, "
        "7 smart filter types, live ticket data, and a world-class 2026 glassmorphism interface."
    )
    add_text_box(slide, summary, Inches(0.5), Inches(2.35), Inches(7.5), Inches(3.2),
                 font_size=13.5, color=RGBColor(0xCC, 0xDD, 0xFF))

    # Key outcomes
    outcomes = [
        ("Natural Language\nBus Search", PURPLE),
        ("Live RedBus\nData", TEAL),
        ("Smart AI\nFiltering", GOLD),
        ("Premium UI\n& Theme", RED_ACCENT),
    ]
    for i, (txt, color) in enumerate(outcomes):
        x = Inches(0.5 + i * 1.85)
        y = Inches(5.7)
        add_rect(slide, x, y, Inches(1.7), Inches(1.0), RGBColor(0x1E, 0x30, 0x5A))
        add_rect(slide, x, y, Inches(1.7), Inches(0.09), color)
        add_text_box(slide, txt, x + Inches(0.08), y + Inches(0.18),
                     Inches(1.55), Inches(0.75), font_size=11, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER)

    # Right side — tech summary
    add_rect(slide, Inches(8.4), Inches(1.55), Inches(4.6), Inches(5.3),
             RGBColor(0x12, 0x1C, 0x3C))
    add_text_box(slide, "Tech at a Glance",
                 Inches(8.6), Inches(1.65), Inches(4.2), Inches(0.4),
                 font_size=13, bold=True, color=TEAL)
    tech_lines = [
        "Frontend:     React + TypeScript + Vite",
        "Styling:      Vanilla CSS Glassmorphism",
        "Backend:      FastAPI (Python 3.11+)",
        "AI:           OpenRouter / Nous Hermes 3",
        "Fallback:     Rule-based Regex Parser",
        "Scraper:      Playwright + RedBus API",
        "Map:          Leaflet.js",
        "Voice:        Web Speech API",
        "Storage:      localStorage",
        "Container:    Docker",
        "Version Ctrl: Git",
    ]
    for i, line in enumerate(tech_lines):
        add_text_box(slide, line, Inches(8.6), Inches(2.15 + i * 0.42),
                     Inches(4.2), Inches(0.38), font_size=11,
                     color=RGBColor(0xCC, 0xDD, 0xFF))

    add_text_box(slide,
        "http://localhost:5173  |  http://localhost:8000/docs",
        Inches(0.5), Inches(6.95), Inches(12), Inches(0.35),
        font_size=10, color=MUTED, align=PP_ALIGN.CENTER
    )


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = prs_init()
    TOTAL = 14

    slide_01_title(prs)
    slide_02_agenda(prs, TOTAL)
    slide_03_overview(prs, TOTAL)
    slide_04_architecture(prs, TOTAL)
    slide_05_tech_frontend(prs, TOTAL)
    slide_06_tech_backend(prs, TOTAL)
    slide_07_ai_pipeline(prs, TOTAL)
    slide_08_filtering(prs, TOTAL)
    slide_09_scraping(prs, TOTAL)
    slide_10_ui_features(prs, TOTAL)
    slide_11_advanced(prs, TOTAL)
    slide_12_dataflow(prs, TOTAL)
    slide_13_limitations(prs, TOTAL)
    slide_14_conclusion(prs, TOTAL)

    prs.save(OUT)
    print(f"PPT saved to: {OUT}")


if __name__ == "__main__":
    main()
